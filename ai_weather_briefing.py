from google import genai
from google.genai import types
import os
import sys
import get_ECMWF_functions as gef
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn

prefix=os.environ["MAIN_PATH"]

if "DATE_STR" in os.environ:
    date_str=os.environ["DATE_STR"]
else:
    today = datetime.today()
    two_days_earlier = today - timedelta(days=2)
    date_str = two_days_earlier.strftime("%Y-%m-%d")

with open(f"{prefix}/prompts/system_prompt.md") as f:
            system_prompt = f.read()

def _load_or_empty(path):
    try:
        return gef.load_dict(path)
    except FileNotFoundError:
        return {}

promt_unformat1=_load_or_empty(f"{prefix}/promt_unformat1.json")
promt_unformat2=_load_or_empty(f"{prefix}/promt_unformat2.json")
promt_unformat3=_load_or_empty(f"{prefix}/promt_unformat3.json")

promt_unformat= promt_unformat1 | promt_unformat2 | promt_unformat3
user_prompt = f"""
Forecast date: {date_str}
Country: Kenya
Month: {date_str[5:7]}
Zone statistics (6-week forecast):
{gef.format_prompt_data(promt_unformat)}
"""

client = genai.Client(api_key=os.environ["GOOGLE_API_KEY"])

response = client.models.generate_content(
    model="gemini-3.1-flash-lite",
    contents=user_prompt,
    config=types.GenerateContentConfig(
        system_instruction=system_prompt,
        max_output_tokens=6000,
    )
)
summary = response.text

# var_ex='''\n \nLegend:\np33= Percentage of ensemble members below normal of model climate
# p66= Percentage of ensemble members above normal of model climate
# p50anom= Anomaly of the ensemble mean from the median of the model climate in % and mm
# efi= Extreme forecast index'''

# with open(f'{prefix}/prompts/digest_{date_str}.txt', 'w') as f:
#     f.write(summary)

def set_slide_text(shape, text, font_size=12, font_name="Calibri", align=None):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    paragraphs = [p.strip() for p in text.strip().split("\n\n") if p.strip()]
    if not paragraphs:
        paragraphs = [text.strip()]

    for idx, para_text in enumerate(paragraphs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()

        if align is not None:
            p.alignment = align

        pPr = p._p.get_or_add_pPr()
        for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
            el = pPr.find(qn(tag))
            if el is not None:
                pPr.remove(el)
        pPr.append(pPr.makeelement(qn("a:buNone"), {}))
        pPr.set("marL", "0")
        pPr.set("indent", "0")

        p.space_after = Pt(8)

        colon_idx = para_text.find(":")
        if 4 < colon_idx < 250:
            header = para_text[:colon_idx + 1]
            rest = para_text[colon_idx + 1:]

            bold_run = p.add_run()
            bold_run.text = header
            bold_run.font.bold = True
            bold_run.font.size = Pt(font_size)
            bold_run.font.name = font_name

            rest_run = p.add_run()
            rest_run.text = rest
            rest_run.font.bold = False
            rest_run.font.size = Pt(font_size)
            rest_run.font.name = font_name
        else:
            run = p.add_run()
            run.text = para_text
            run.font.size = Pt(font_size)
            run.font.name = font_name

prs = Presentation("WeatherbriefingKenya_template3.pptx")
text = summary
slide_text = text.split("---SLIDE---")

slide_types = ["date", 'ECMWF_raw', 'GEFS_raw', "ECMWF_dwn", "ECMWFmed_raw",
               "ECMWF_tercile", "ECMWF_efi", "ECMWF_p50", "ECMWF_p50anom", "sum", "waves", "IOD"]
plots = ['hold', 'weekly_precip', "gefs_weekly_precip", "weekly_precip_downscaled",
         "weekly_medium_range_precip", "chance_of_above_or_below", "efi_sot_precip",
         "50th_percentile_exedance", "anomaly_from_50th", "summary"]

kenya_path = f"plots/Kenya/{date_str}"
great_horn_path = f"plots/Great_Horn/{date_str}"
diagnostics_path = f"plots/diagnostics/{date_str}/weekly"
diagnostics_monthly_path = f"plots/diagnostics/{date_str}/monthly"
briefing_plots_path = f"plots/briefing/{date_str}/"

# only the first len(plots) slide types have a "{type}_plot" shape
plot_paths = {t: f"{kenya_path}/weekly/{p}.png" for t, p in zip(slide_types, plots)}

# Indian Ocean moisture diagnostics (see IndianOceanState.py)
IOD_path = f"{diagnostics_path}/ECMWF_s2s_10wind_sst_anomaly_{date_str}.png"
IO_ivt_weekly_path = f"{diagnostics_path}/ECMWF_s2s_ivt_u_{date_str}.png"
IO_ivt_monthly_path = f"{diagnostics_monthly_path}/ECMWF_s2s_ivt_u_{date_str}.png"
IO_TCWV_anom_path = f"{diagnostics_path}/ECMWF_s2s_tcw_anomaly_{date_str}.png"
IO_precip_anom_path = f"{diagnostics_path}/ECMWF_s2s_precip_anomaly_{date_str}.png"
IO_precip_anom_std_path = f"{diagnostics_path}/ECMWF_s2s_precip_std_anomaly_{date_str}.png"

# rainy season onset maps (see run_rainfall_onset.py) -- wet-spell/no-dry-spell
# definition, and the two-stage cumulative-rainfall ("accum") definition
onsetecmwf_path = f"{kenya_path}/monthly/onset_s2s.png"
onsetgefs_path = f"{kenya_path}/monthly/onset_gefs.png"
onsetecmwf_accum_path = f"{kenya_path}/monthly/onset_s2s_accum.png"
onsetgefs_accum_path = f"{kenya_path}/monthly/onset_gefs_accum.png"

# dry/wet spell probability & median length maps (see plot_s2s.py)
median_wet_path = f"{kenya_path}/monthly/median_wetspell_length.png"
wet5_path = f"{kenya_path}/monthly/prob_wetspell_5days.png"
wet7_path = f"{kenya_path}/monthly/prob_wetspell_7days.png"
median_dry_path = f"{kenya_path}/monthly/median_dryspell_length.png"
dry5_path = f"{kenya_path}/monthly/prob_dryspell_5days.png"
dry7_path = f"{kenya_path}/monthly/prob_dryspell_7days.png"

exceed20mm_path = f"{kenya_path}/weekly/chance_higherthan_20mm.png"

# Great Horn tercile plot (regional counterpart to the Kenya-only ECMWF_tercile_plot)
ecmwf_tercile_ea_path = f"{great_horn_path}/weekly/chance_of_above_or_below.png"

picture_paths = {
    "ECMWF_tercile_plot_EA": ecmwf_tercile_ea_path,
    "IO_state": IOD_path,
    "IO_ivt_weekly": IO_ivt_weekly_path,
    "IO_ivt_monthly": IO_ivt_monthly_path,
    "IO_TCWV_anom": IO_TCWV_anom_path,
    "IO_precip_anom": IO_precip_anom_path,
    "IO_precip_anom_std": IO_precip_anom_std_path,
    "Onset_ECMWF": onsetecmwf_path,
    "Onset_GEFS": onsetgefs_path,
    "Onset_ECMWF_accum": onsetecmwf_accum_path,
    "Onset_GEFS_accum": onsetgefs_accum_path,
    "median_wet": median_wet_path,
    "wet5": wet5_path,
    "wet7": wet7_path,
    "median_dry": median_dry_path,
    "dry5": dry5_path,
    "dry7": dry7_path,
    "exceed20mm": exceed20mm_path,
}

# Plots generated by ws_scripts/slide*.sh and collected into briefing_plots_path
# (see populate_briefing_template3.py) — picture-only, no AI narration.
briefing_plot_names = [
    "chirps_kenya_weekly_rainfall",
    "chirps_kenya_weekly_anomaly",
    "chirps_east_africa_weekly_rainfall",
    "chirps_east_africa_weekly_anomaly",
    "tahmo_kenya_cities_rainfall",
    "tahmo_kenya_cities_temperature",
    "tahmo_kenya_cities_humidity",
    "kenya_ond_weekly_rainfall_vs_climatology",
    "kenya_ond_weekly_standardized_anomaly",
    "kenya_weekly_rainfall_analog_years",
    "mjo_rmm_gefs",
    "mjo_rmm_ecmwf",
    "iod_observed",
    "enso_observed",
    "iod_forecast",
    "enso_forecast",
    "itcz_africa_latest",
    "hovmoller_olr_tropics",
    "olr_map_africa",
    "kenya_gefs_chirps_verify_5mm",
    "kenya_gefs_chirps_bias",
    "kenya_gefs_chirps_mae",
]
for name in briefing_plot_names:
    picture_paths[name] = f"{briefing_plots_path}/{name}.png"

def replace_picture(slide, shape, image_path):
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    shape._element.getparent().remove(shape._element)
    slide.shapes.add_picture(image_path, left, top, width, height)

# Picture-only shapes (no AI narration) — matched by exact shape name,
# wherever in the deck that shape happens to live. Skip (warn, don't crash)
# when the source file isn't there yet -- e.g. plots/{date}/briefing_plots
# is only populated by ws_scripts/slide*.sh today, not yet by this pipeline.
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.name in picture_paths:
            path = picture_paths[shape.name]
            if os.path.exists(path):
                replace_picture(slide, shape, path)
            else:
                print(f"WARNING: missing picture for '{shape.name}': {path}", file=sys.stderr)

dt_obj = datetime.fromisoformat(date_str)
day = dt_obj.day
suffix = ("th" if 11 <= day <= 13 else {1: "st", 2: "nd", 3: "rd"}.get(day % 10, "th"))
formatted_date = f"{dt_obj.strftime('%B')} {day}{suffix}, {dt_obj.strftime('%Y')}"

month_abbrevs = {1: "Jan", 2: "Feb", 3: "Mar", 4: "Apr", 5: "May", 6: "Jun",
                  7: "Jul", 8: "Aug", 9: "Sept", 10: "Oct", 11: "Nov", 12: "Dec"}
short_date = f"{day} {month_abbrevs[dt_obj.month]} {dt_obj.strftime('%Y')}"

# AI-narrated slide types. Template3 no longer places "{type}_text"/
# "{type}_plot" shapes at the same slide index as slide_types (they're
# scattered among many new picture-only slides), so match by name across
# the whole deck instead of assuming slide i holds slide_types[i].
for t, text in zip(slide_types, slide_text):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name == f"{t}_text":
                if t == 'date':
                    set_slide_text(shape, formatted_date, font_size=35, font_name='Karla Medium', align=PP_ALIGN.CENTER)
                elif t == 'sum':
                    set_slide_text(shape, text, font_size=17)
                else:
                    set_slide_text(shape, text, font_size=13)
            elif shape.name == f"{t}_plot":
                replace_picture(slide, shape, plot_paths[t])

# Extra date-bearing shapes that don't follow the "{type}_text" naming
# convention: "gen_date" has a "-date-" placeholder inline in a longer
# sentence, and "sum title" is a short "<day> <Mon> <year> Outlook" heading.
# Both are edited in-place (existing run text) to preserve their template
# formatting rather than going through set_slide_text.
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.name == "gen_date":
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    if "-date-" in run.text:
                        run.text = run.text.replace("-date-", formatted_date)
        elif shape.name == "sum title":
            runs = shape.text_frame.paragraphs[0].runs
            if runs:
                runs[0].text = f"{short_date} Outlook"
                for run in runs[1:]:
                    run.text = ""

prs.save(f"s2s_briefing_{date_str}.pptx")